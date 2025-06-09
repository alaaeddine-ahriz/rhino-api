"""Document processing and text splitting for RAG system."""
import os
import hashlib
import json
import shutil
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import glob

from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter
from langchain.schema import Document

from app.core.config import settings
from app.services.rag.core import initialize_pinecone, setup_embeddings, create_or_get_index
from app.services.rag.embeddings import upsert_documents

def initialiser_structure_dossiers():
    """
    Initialize the folder structure for courses if it doesn't exist.
    Creates a folder per subject with an explanatory README file.
    """
    # Create main courses folder if it doesn't exist
    if not os.path.exists(settings.COURS_DIR):
        os.makedirs(settings.COURS_DIR)
        print(f"Main courses folder created: {settings.COURS_DIR}")
    
    # List of subjects (to be adapted according to needs)
    matieres = ["SYD", "TCP"]
    
    for matiere in matieres:
        matiere_dir = os.path.join(settings.COURS_DIR, matiere)
        if not os.path.exists(matiere_dir):
            os.makedirs(matiere_dir)
            
            # Create an explanatory README file
            readme_path = os.path.join(matiere_dir, "README.md")
            with open(readme_path, "w") as f:
                f.write(f"# Course documents for {matiere}\n\n")
                f.write("Place your course documents for this subject here.\n")
                f.write("Supported formats: .md, .txt\n\n")
                f.write("Recommended structure for markdown files:\n")
                f.write("- Use ## for main sections\n")
                f.write("- Each file should cover one concept or topic\n")
            
            print(f"Folder for subject {matiere} created with explanatory README")

def lire_fichiers_matiere(matiere: str) -> List[Dict[str, Any]]:
    """
    Read all course files for a given subject,
    with support for various formats (md, txt, pdf, docx, pptx, etc.).
    
    Args:
        matiere: Subject identifier
        
    Returns:
        List[Dict[str, Any]]: List of documents with their content and metadata
    """
    matiere_dir = os.path.join(settings.COURS_DIR, matiere)
    
    # Check if folder exists
    if not os.path.exists(matiere_dir):
        print(f"Error: Subject folder {matiere} does not exist.")
        return []
    
    documents = []
    
    # Supported extensions
    extensions = ["*.md", "*.txt", "*.pdf", "*.docx", "*.pptx", "*.doc", "*.odt", "*.odp"]
    
    # Check if an exams folder exists for this subject
    examens_dir = os.path.join(matiere_dir, "examens")
    exam_documents = []
    has_exam_folder = os.path.exists(examens_dir) and os.path.isdir(examens_dir)
    
    # Process all files with supported extensions
    for ext in extensions:
        for file_path in glob.glob(os.path.join(matiere_dir, "**", ext), recursive=True):
            try:
                # Skip README files
                if os.path.basename(file_path).lower() == "readme.md":
                    continue
                
                # Calculate file hash to track modifications
                file_hash = calculer_hash_fichier(file_path)
                
                # Extract content based on file type
                file_extension = os.path.splitext(file_path)[1].lower()
                content = extraire_contenu_fichier(file_path, file_extension)
                
                if not content or content.strip() == "":
                    print(f"Warning: File {file_path} seems empty after extraction.")
                    continue
                
                # Document metadata
                relative_path = os.path.relpath(file_path, settings.COURS_DIR)
                metadata = {
                    "source": relative_path,
                    "matiere": matiere,
                    "filename": os.path.basename(file_path),
                    "filetype": file_extension,
                    "file_hash": file_hash,
                    "updated_at": datetime.now().isoformat()
                }
                
                # Check if document is in exams folder
                is_exam = "examens" in relative_path
                if is_exam:
                    metadata["is_exam"] = True
                    metadata["document_type"] = "exam"
                    exam_documents.append({"content": content, "metadata": metadata})
                else:
                    documents.append({"content": content, "metadata": metadata})
                
                print(f"File read: {relative_path}" + (" (exam)" if is_exam else ""))
                
            except Exception as e:
                print(f"Error reading file {file_path}: {e}")
    
    # Combine documents, placing exams first to give them more weight
    return exam_documents + documents

def calculer_hash_fichier(file_path: str) -> str:
    """
    Calculate an MD5 hash of a file's content to detect modifications.
    
    Args:
        file_path: Path to the file
        
    Returns:
        str: MD5 hash of the file
    """
    hash_md5 = hashlib.md5()
    
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
            
    return hash_md5.hexdigest()

def extraire_contenu_fichier(file_path: str, file_extension: str) -> str:
    """
    Extract text content from a file based on its format.
    
    Args:
        file_path: Path to the file
        file_extension: File extension (with dot)
        
    Returns:
        str: Text content of the file
    """
    try:
        # Text and markdown files
        if file_extension in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # PDF files
        elif file_extension == '.pdf':
            try:
                import pdfplumber
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        text += page_text + "\n\n"
                return text
            except Exception as e:
                print(f"Error with pdfplumber: {e}. Trying PyPDF2...")
                
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as file:
                        reader = PyPDF2.PdfReader(file)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text() + "\n\n"
                    return text
                except Exception as pdf_err:
                    print(f"Error with PyPDF2: {pdf_err}")
                    return f"[PDF extraction error: {str(pdf_err)}]"
        
        # Word files (DOCX)
        elif file_extension == '.docx':
            try:
                import docx
                text = ""
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
                    # Try to identify potential headers based on style
                    if para.style.name.startswith('Heading'):
                        heading_level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 2
                        prefix = '#' * heading_level
                        text = text[:-1] + f"\n{prefix} {para.text}\n"
                return text
            except Exception as e:
                print(f"Error reading DOCX file: {e}")
                return f"[DOCX extraction error: {str(e)}]"
        
        # PowerPoint files (PPTX)
        elif file_extension == '.pptx':
            try:
                from pptx import Presentation
                text = ""
                pres = Presentation(file_path)
                for i, slide in enumerate(pres.slides):
                    text += f"## Slide {i+1}\n\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                    text += "\n"
                return text
            except Exception as e:
                print(f"Error reading PPTX file: {e}")
                return f"[PPTX extraction error: {str(e)}]"
        
        # OpenDocument Text (ODT)
        elif file_extension == '.odt':
            try:
                import odf.opendocument
                from odf.text import P
                from odf.teletype import extractText
                
                textdoc = odf.opendocument.load(file_path)
                allparas = textdoc.getElementsByType(P)
                text = ""
                for para in allparas:
                    text += extractText(para) + "\n"
                return text
            except Exception as e:
                print(f"Error reading ODT file: {e}")
                return f"[ODT extraction error: {str(e)}]"
        
        # OpenDocument Presentation (ODP)
        elif file_extension == '.odp':
            try:
                import odf.opendocument
                from odf.text import P
                from odf.teletype import extractText
                
                doc = odf.opendocument.load(file_path)
                text = ""
                slide_num = 1
                
                # Get all text elements
                for para in doc.getElementsByType(P):
                    content = extractText(para)
                    if content.strip():
                        if "Slide" not in text[-20:] and len(text) > 0:
                            text += f"\n## Slide {slide_num}\n\n"
                            slide_num += 1
                        text += content + "\n"
                
                return text
            except Exception as e:
                print(f"Error reading ODP file: {e}")
                return f"[ODP extraction error: {str(e)}]"
        
        # DOC files (old Word format)
        elif file_extension == '.doc':
            try:
                import textract
                text = textract.process(file_path).decode('utf-8')
                return text
            except ImportError:
                print(f"Error: textract not installed for reading .doc files")
                print(f"Install it with 'pip install textract'")
                return f"[DOC content not extracted: {os.path.basename(file_path)}]"
            except Exception as e:
                print(f"Error reading DOC file: {e}")
                return f"[DOC extraction error: {str(e)}]"
        
        # Unsupported format
        else:
            print(f"Unsupported file format: {file_extension}")
            return f"[Unsupported format: {file_extension}]"
    
    except Exception as e:
        print(f"Error extracting content from {file_path}: {e}")
        return f"[Extraction error: {str(e)}]"

def split_document(document: Dict[str, Any]) -> List[Document]:
    """
    Split a document into sections.
    Uses different strategies based on file type and header presence.
    
    Args:
        document: Document with content and metadata
        
    Returns:
        List[Document]: List of sections with their metadata
    """
    content = document["content"]
    metadata = document["metadata"]
    
    # Check for markdown headers (##, ###) in content
    has_markdown_headers = '##' in content
    
    # If document is markdown or contains headers, use header-based method
    if metadata["filetype"] == ".md" or has_markdown_headers:
        headers_to_split_on = [
            ("##", "Header 2"),
            ("###", "Header 3")
        ]
        
        try:
            markdown_splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=headers_to_split_on, strip_headers=False
            )
            
            splits = markdown_splitter.split_text(content)
            
            # If no headers found, use character-based method
            if not splits:
                return split_by_characters(content, metadata)
                
            # Add document metadata to each split
            for split in splits:
                split.metadata.update(metadata)
                
            return splits
            
        except Exception as e:
            print(f"Error during markdown splitting: {e}")
            # Fallback to character-based splitting
            return split_by_characters(content, metadata)
    
    # For other file types or documents without headers: split by characters and paragraphs
    else:
        return split_by_paragraphs(content, metadata)

def split_by_characters(content: str, metadata: Dict[str, Any]) -> List[Document]:
    """
    Split a document into chunks by characters (fallback method).
    
    Args:
        content: Document content
        metadata: Document metadata
        
    Returns:
        List[Document]: List of sections with their metadata
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    
    docs = text_splitter.create_documents([content], [metadata])
    return docs

def split_by_paragraphs(content: str, metadata: Dict[str, Any]) -> List[Document]:
    """
    Split a document into chunks by paragraphs,
    more suitable for unstructured documents.
    
    Args:
        content: Document content
        metadata: Document metadata
        
    Returns:
        List[Document]: List of sections with their metadata
    """
    # Split by paragraphs with overlap
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=250,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )
    
    docs = text_splitter.create_documents([content], [metadata])
    return docs

def get_documents_for_subject(matiere: str) -> List[Dict[str, Any]]:
    """
    Get list of documents for a specific subject with metadata.
    
    Args:
        matiere: Subject identifier
        
    Returns:
        List[Dict[str, Any]]: List of documents with their metadata
    """
    matiere_dir = os.path.join(settings.COURS_DIR, matiere)
    
    if not os.path.exists(matiere_dir):
        return []
    
    documents = []
    extensions = ["*.md", "*.txt", "*.pdf", "*.docx", "*.pptx", "*.doc", "*.odt", "*.odp"]
    
    for ext in extensions:
        for file_path in glob.glob(os.path.join(matiere_dir, "**", ext), recursive=True):
            # Skip README files
            if os.path.basename(file_path).lower() == "readme.md":
                continue
                
            try:
                relative_path = os.path.relpath(file_path, settings.COURS_DIR)
                file_stats = os.stat(file_path)
                file_hash = calculer_hash_fichier(file_path)
                
                # Check if document is in exams folder
                is_exam = "examens" in relative_path
                
                document_info = {
                    "id": file_hash,  # Using file hash as unique ID
                    "filename": os.path.basename(file_path),
                    "matiere": matiere,
                    "document_type": os.path.splitext(file_path)[1].lower().lstrip('.'),
                    "is_exam": is_exam,
                    "file_path": relative_path,
                    "file_size": file_stats.st_size,
                    "upload_date": datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
                    "last_modified": datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
                    "file_hash": file_hash
                }
                
                documents.append(document_info)
                
            except Exception as e:
                print(f"Error getting info for file {file_path}: {e}")
    
    # Sort by upload date (newest first)
    documents.sort(key=lambda x: x["upload_date"], reverse=True)
    return documents

def upload_document_to_subject(
    matiere: str, 
    filename: str, 
    file_content: bytes, 
    is_exam: bool = False
) -> Tuple[bool, str, Optional[Dict[str, Any]]]:
    """
    Upload a document to a specific subject folder.
    
    Args:
        matiere: Subject identifier
        filename: Name of the file
        file_content: File content as bytes
        is_exam: Whether this is an exam document
        
    Returns:
        Tuple[bool, str, Optional[Dict]]: (success, message, document_info)
    """
    try:
        # Ensure subject folder exists
        matiere_dir = os.path.join(settings.COURS_DIR, matiere)
        if not os.path.exists(matiere_dir):
            os.makedirs(matiere_dir)
        
        # If it's an exam, create/use exams subfolder
        if is_exam:
            target_dir = os.path.join(matiere_dir, "examens")
            if not os.path.exists(target_dir):
                os.makedirs(target_dir)
        else:
            target_dir = matiere_dir
        
        # Validate file extension
        file_extension = os.path.splitext(filename)[1].lower()
        allowed_extensions = ['.md', '.txt', '.pdf', '.docx', '.pptx', '.doc', '.odt', '.odp']
        
        if file_extension not in allowed_extensions:
            return False, f"File type {file_extension} not supported. Allowed types: {', '.join(allowed_extensions)}", None
        
        # Create full file path
        file_path = os.path.join(target_dir, filename)
        
        # Check if file already exists
        if os.path.exists(file_path):
            return False, f"File {filename} already exists in {matiere}", None
        
        # Write file content
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        # Get file info
        file_stats = os.stat(file_path)
        file_hash = calculer_hash_fichier(file_path)
        relative_path = os.path.relpath(file_path, settings.COURS_DIR)
        
        document_info = {
            "id": file_hash,
            "filename": filename,
            "matiere": matiere,
            "document_type": file_extension.lstrip('.'),
            "is_exam": is_exam,
            "file_path": relative_path,
            "file_size": file_stats.st_size,
            "upload_date": datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
            "last_modified": datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
            "file_hash": file_hash
        }
        
        return True, f"Document {filename} uploaded successfully", document_info
        
    except Exception as e:
        return False, f"Error uploading document: {str(e)}", None

def delete_document_from_subject(matiere: str, document_id: str) -> Tuple[bool, str]:
    """
    Delete a document from a subject folder.
    
    Args:
        matiere: Subject identifier
        document_id: Document ID (file hash)
        
    Returns:
        Tuple[bool, str]: (success, message)
    """
    try:
        # Find the document by ID (hash)
        documents = get_documents_for_subject(matiere)
        target_document = None
        
        for doc in documents:
            if doc["id"] == document_id:
                target_document = doc
                break
        
        if not target_document:
            return False, f"Document with ID {document_id} not found in {matiere}"
        
        # Construct full file path
        full_path = os.path.join(settings.COURS_DIR, target_document["file_path"])
        
        if not os.path.exists(full_path):
            return False, f"File {target_document['filename']} not found on disk"
        
        # Delete the file
        os.remove(full_path)
        
        return True, f"Document {target_document['filename']} deleted successfully"
        
    except Exception as e:
        return False, f"Error deleting document: {str(e)}"

def get_document_content(matiere: str, document_id: str) -> Tuple[bool, str, Optional[str]]:
    """
    Get the content of a specific document.
    
    Args:
        matiere: Subject identifier
        document_id: Document ID (file hash)
        
    Returns:
        Tuple[bool, str, Optional[str]]: (success, message, content)
    """
    try:
        # Find the document by ID
        documents = get_documents_for_subject(matiere)
        target_document = None
        
        for doc in documents:
            if doc["id"] == document_id:
                target_document = doc
                break
        
        if not target_document:
            return False, f"Document with ID {document_id} not found", None
        
        # Get full file path
        full_path = os.path.join(settings.COURS_DIR, target_document["file_path"])
        
        if not os.path.exists(full_path):
            return False, f"File not found on disk", None
        
        # Extract content based on file type
        content = extraire_contenu_fichier(full_path, target_document["document_type"])
        
        if not content:
            return False, f"Could not extract content from file", None
        
        return True, "Content extracted successfully", content
        
    except Exception as e:
        return False, f"Error extracting content: {str(e)}", None

def process_and_index_new_document(
    matiere: str, 
    document_info: Dict[str, Any]
) -> Tuple[bool, str]:
    """
    Process and index a newly uploaded document into the vector database.
    
    Args:
        matiere: Subject identifier
        document_info: Document information returned from upload_document_to_subject
        
    Returns:
        Tuple[bool, str]: (success, message)
    """
    try:
        # Initialize RAG components
        pc, index_name, spec = initialize_pinecone()
        embeddings = setup_embeddings()
        vector_store = create_or_get_index(pc, index_name, embeddings, spec)
        
        # Create document object for processing
        full_path = os.path.join(settings.COURS_DIR, document_info["file_path"])
        content = extraire_contenu_fichier(full_path, f".{document_info['document_type']}")
        
        if not content:
            return False, f"Could not extract content from uploaded document"
        
        # Create document with metadata
        document = {
            "content": content,
            "metadata": {
                "source": document_info["file_path"],
                "matiere": matiere,
                "filename": document_info["filename"],
                "filetype": f".{document_info['document_type']}",
                "file_hash": document_info["file_hash"],
                "is_exam": document_info["is_exam"],
                "document_type": "exam" if document_info["is_exam"] else "course",
                "updated_at": document_info["upload_date"]
            }
        }
        
        # Split document into sections
        sections = split_document(document)
        
        if not sections:
            return False, f"Document could not be split into sections"
        
        # Index the document sections
        vector_store_result, namespace = upsert_documents(
            pc=pc,
            index_name=index_name,
            embeddings=embeddings,
            matiere=matiere,
            docs=sections
        )
        
        if vector_store_result:
            return True, f"Document successfully indexed with {len(sections)} sections in namespace '{namespace}'"
        else:
            return False, f"Failed to index document sections"
            
    except Exception as e:
        return False, f"Error processing and indexing document: {str(e)}"

def reindex_document_if_modified(
    matiere: str, 
    document_info: Dict[str, Any]
) -> Tuple[bool, str]:
    """
    Check if a document needs reindexing and perform it if necessary.
    
    Args:
        matiere: Subject identifier
        document_info: Document information
        
    Returns:
        Tuple[bool, str]: (success, message)
    """
    try:
        # This function can be used for future automatic reindexing
        # when documents are detected as modified
        return process_and_index_new_document(matiere, document_info)
        
    except Exception as e:
        return False, f"Error reindexing document: {str(e)}"

if __name__ == "__main__":
    matiere = "SYD"
    documents = lire_fichiers_matiere(matiere)
    print(documents)