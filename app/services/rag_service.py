"""
Service RAG autonome - Système de Retrieval-Augmented Generation
Basé sur l'analyse du projet rhino-v1-deprecated

Ce module fournit une classe RAGService qui encapsule toute la fonctionnalité
de génération augmentée par récupération (RAG) dans un seul fichier réutilisable.

Fonctionnalités:
- Gestion de documents multi-formats (PDF, DOCX, TXT, MD)
- Base de données vectorielle Pinecone avec namespaces
- Génération de questions et évaluation de réponses
- Support des examens avec métadonnées spéciales
- Prompts personnalisés pour différents cas d'usage

Usage:
    rag = RAGService(
        pinecone_api_key="your_key",
        openai_api_key="your_key",
        index_name="your_index"
    )
    
    # Ajouter des documents
    rag.add_documents("matiere1", documents)
    
    # Poser une question
    response = rag.query("matiere1", "Qu'est-ce que...?")
    
    # Générer une question de réflexion
    question = rag.generate_reflection_question("matiere1", "concept_cle")
"""

import os
import time
import uuid
import warnings
import glob
import hashlib
import json
import tempfile
import io
import re
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Union

# Désactivation des avertissements LangSmith
os.environ["LANGCHAIN_TRACING_V2"] = "false"
os.environ["LANGCHAIN_TRACING"] = "false"
warnings.filterwarnings("ignore", category=Warning, module="langsmith")

# Imports pour les différents formats de documents
try:
    import pdfplumber
    import PyPDF2
    import docx
    from pptx import Presentation
    import odf.opendocument
    from odf.text import P
    from odf.teletype import extractText
    PDF_SUPPORT = True
except ImportError as e:
    print(f"Avertissement: Certaines bibliothèques de traitement de documents ne sont pas installées: {e}")
    PDF_SUPPORT = False

# Imports pour le découpage de texte
from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter

# Imports pour la base de données vectorielle
from pinecone import Pinecone, ServerlessSpec
from langchain_pinecone import PineconeEmbeddings, PineconeVectorStore

# Imports pour le système RAG
from langchain_openai import ChatOpenAI
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.prompts import ChatPromptTemplate
from langchain import hub
from langchain.schema import Document


class RAGService:
    """
    Service RAG autonome pour la génération augmentée par récupération.
    
    Cette classe encapsule toute la fonctionnalité RAG nécessaire pour:
    - Gérer des documents de cours par matière
    - Effectuer des recherches sémantiques
    - Générer des réponses contextuelles
    - Créer des questions de réflexion
    - Évaluer des réponses d'étudiants
    """
    
    def __init__(
        self,
        pinecone_api_key: str,
        openai_api_key: str,
        index_name: str = "rag-index",
        embedding_model: str = "multilingual-e5-large",
        llm_model: str = "gpt-4o-mini",
        pinecone_cloud: str = "aws",
        pinecone_region: str = "us-east-1"
    ):
        """
        Initialise le service RAG.
        
        Args:
            pinecone_api_key: Clé API Pinecone
            openai_api_key: Clé API OpenAI
            index_name: Nom de l'index Pinecone
            embedding_model: Modèle d'embedding à utiliser
            llm_model: Modèle de langage à utiliser
            pinecone_cloud: Cloud provider pour Pinecone
            pinecone_region: Région pour Pinecone
        """
        self.pinecone_api_key = pinecone_api_key
        self.openai_api_key = openai_api_key
        self.index_name = index_name
        self.embedding_model = embedding_model
        self.llm_model = llm_model
        
        # Configuration des variables d'environnement
        os.environ["PINECONE_API_KEY"] = pinecone_api_key
        os.environ["OPENAI_API_KEY"] = openai_api_key
        
        # Initialisation des composants
        self.pc = Pinecone(api_key=pinecone_api_key)
        self.spec = ServerlessSpec(cloud=pinecone_cloud, region=pinecone_region)
        self.embeddings = PineconeEmbeddings(
            model=embedding_model,
            pinecone_api_key=pinecone_api_key
        )
        
        # Créer ou récupérer l'index
        self.index = self._create_or_get_index()
        
        # Cache pour les métadonnées des documents
        self.metadata_cache = {}
    
    def _create_or_get_index(self):
        """Crée un nouvel index si nécessaire ou récupère un index existant."""
        if self.index_name not in self.pc.list_indexes().names():
            print(f"Création d'un nouvel index: {self.index_name}")
            self.pc.create_index(
                name=self.index_name,
                dimension=self.embeddings.dimension,
                metric="cosine",
                spec=self.spec
            )
            # Attendre que l'index s'initialise
            print("Attente de l'initialisation de l'index...")
            time.sleep(10)
        
        return self.pc.Index(self.index_name)
    
    def _get_namespace(self, matiere: str) -> str:
        """Génère un namespace standardisé pour une matière."""
        return f"matiere-{matiere.lower()}"
    
    def _calculate_file_hash(self, content: str) -> str:
        """Calcule un hash MD5 du contenu pour détecter les modifications."""
        return hashlib.md5(content.encode('utf-8')).hexdigest()
    
    def _extract_content_from_file(self, file_path: str) -> str:
        """
        Extrait le contenu textuel d'un fichier selon son extension.
        
        Args:
            file_path: Chemin vers le fichier
            
        Returns:
            Contenu textuel du fichier
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            
            elif file_extension == '.md':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            
            elif file_extension == '.pdf' and PDF_SUPPORT:
                return self._extract_pdf_content(file_path)
            
            elif file_extension == '.docx' and PDF_SUPPORT:
                doc = docx.Document(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
            elif file_extension == '.pptx' and PDF_SUPPORT:
                prs = Presentation(file_path)
                text_content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
                return '\n'.join(text_content)
            
            else:
                print(f"Format de fichier non supporté: {file_extension}")
                return ""
                
        except Exception as e:
            print(f"Erreur lors de l'extraction du contenu de {file_path}: {e}")
            return ""
    
    def _extract_pdf_content(self, file_path: str) -> str:
        """Extrait le contenu d'un fichier PDF."""
        try:
            # Essayer avec pdfplumber d'abord
            with pdfplumber.open(file_path) as pdf:
                text_content = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
                return '\n'.join(text_content)
        except:
            try:
                # Fallback avec PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text_content = []
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
                    return '\n'.join(text_content)
            except Exception as e:
                print(f"Erreur lors de l'extraction PDF: {e}")
                return ""
    
    def _split_document(self, content: str, metadata: Dict[str, Any]) -> List[Document]:
        """
        Découpe un document en sections plus petites pour l'indexation.
        
        Args:
            content: Contenu du document
            metadata: Métadonnées du document
            
        Returns:
            Liste de documents découpés
        """
        documents = []
        
        # Si c'est un fichier Markdown, utiliser le découpage par headers
        if metadata.get('filetype') == '.md':
            headers_to_split_on = [
                ("#", "Header 1"),
                ("##", "Header 2"),
                ("###", "Header 3"),
            ]
            
            markdown_splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=headers_to_split_on
            )
            
            try:
                md_header_splits = markdown_splitter.split_text(content)
                for split in md_header_splits:
                    # Combiner les métadonnées
                    combined_metadata = {**metadata, **split.metadata}
                    documents.append(Document(
                        page_content=split.page_content,
                        metadata=combined_metadata
                    ))
                return documents
            except Exception as e:
                print(f"Erreur lors du découpage Markdown: {e}")
        
        # Découpage par caractères pour les autres formats
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
        splits = text_splitter.split_text(content)
        for i, split in enumerate(splits):
            chunk_metadata = metadata.copy()
            chunk_metadata['chunk_id'] = i
            documents.append(Document(
                page_content=split,
                metadata=chunk_metadata
            ))
        
        return documents
    
    def add_documents_from_directory(
        self, 
        matiere: str, 
        directory_path: str,
        supported_extensions: List[str] = None
    ) -> bool:
        """
        Ajoute tous les documents d'un répertoire à l'index.
        
        Args:
            matiere: Identifiant de la matière
            directory_path: Chemin vers le répertoire contenant les documents
            supported_extensions: Extensions de fichiers supportées
            
        Returns:
            True si des documents ont été ajoutés, False sinon
        """
        if supported_extensions is None:
            supported_extensions = ['.md', '.txt', '.pdf', '.docx', '.pptx']
        
        documents = []
        
        # Parcourir récursivement le répertoire
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                file_extension = os.path.splitext(file)[1].lower()
                
                if file_extension in supported_extensions:
                    # Éviter les fichiers README
                    if file.lower() == "readme.md":
                        continue
                    
                    content = self._extract_content_from_file(file_path)
                    if not content or content.strip() == "":
                        continue
                    
                    # Métadonnées du document
                    relative_path = os.path.relpath(file_path, directory_path)
                    metadata = {
                        "source": relative_path,
                        "matiere": matiere,
                        "filename": file,
                        "filetype": file_extension,
                        "file_hash": self._calculate_file_hash(content),
                        "updated_at": datetime.now().isoformat(),
                        "is_exam": "exam" in relative_path.lower()
                    }
                    
                    # Découper le document
                    doc_splits = self._split_document(content, metadata)
                    documents.extend(doc_splits)
                    
                    print(f"Document ajouté: {relative_path}")
        
        if documents:
            return self.add_documents(matiere, documents)
        
        return False
    
    def add_documents(self, matiere: str, documents: List[Document]) -> bool:
        """
        Ajoute des documents à l'index vectoriel.
        
        Args:
            matiere: Identifiant de la matière
            documents: Liste de documents à ajouter
            
        Returns:
            True si l'ajout a réussi, False sinon
        """
        if not documents:
            return False
        
        namespace = self._get_namespace(matiere)
        
        try:
            print(f"Ajout de {len(documents)} documents dans l'espace '{namespace}'")
            vector_store = PineconeVectorStore.from_documents(
                documents=documents,
                index_name=self.index_name,
                embedding=self.embeddings,
                namespace=namespace
            )
            print("Ajout réussi")
            return True
        except Exception as e:
            print(f"Erreur lors de l'ajout: {e}")
            return False
    
    def _create_json_prompt(self, matiere: str) -> ChatPromptTemplate:
        """Crée un prompt pour générer des réponses au format JSON."""
        template = """
        Vous êtes un assistant pédagogique IA spécialisé dans la matière {matiere}, disposant d'un accès direct aux documents de cours via un système de recherche sémantique (RAG).

        Sur la base des extraits de cours suivants:
        {context}
        
        Répondez à cette question: {input}
        
        Votre réponse DOIT être au format JSON suivant:
        
        ```json
        {{
            "réponse": "La réponse complète à la question",
            "niveau_confiance": 0.95,
            "concepts_clés": ["concept1", "concept2", "concept3"]
        }}
        ```
        
        IMPORTANT:
        - Ne mentionnez pas et n'ajoutez pas de sources dans votre réponse
        - Concentrez-vous principalement sur la réponse
        - Le champ "concepts_clés" est optionnel
        
        Ne répondez qu'avec ce format JSON, sans aucun texte avant ou après.
        """
        
        return ChatPromptTemplate.from_template(template).partial(matiere=matiere)
    
    def _create_tutor_prompt(self, matiere: str, output_format: str = "text") -> ChatPromptTemplate:
        """Crée un prompt pour générer des questions de réflexion."""
        if output_format == "json":
            template = """
            Vous êtes un tuteur IA spécialisé dans la matière {matiere}.

            Votre tâche est de générer une question de réflexion originale sur le concept demandé, en vous basant strictement sur les extraits suivants:

            {context}

            La question doit :
            1. Être ORIGINALE et non une reformulation des questions existantes
            2. Solliciter l'analyse critique d'un concept ou d'une relation entre plusieurs notions
            3. Être formulée de manière claire, concise et précise
            4. Favoriser une réponse argumentée plutôt qu'une simple définition

            Votre réponse DOIT être au format JSON suivant:

            ```json
            {{
                "question": "La question de réflexion originale",
                "concepts_abordés": ["concept1", "concept2", "concept3"],
                "niveau_difficulté": "avancé",
                "compétences_visées": ["analyse critique", "synthèse", "application pratique"],
                "éléments_réponse": [
                    "Élément 1 attendu dans la réponse",
                    "Élément 2 attendu dans la réponse"
                ]
            }}
            ```

            Ne répondez qu'avec ce format JSON, sans aucun texte avant ou après.
            """
        else:
            template = """
            Vous êtes un tuteur IA spécialisé dans la matière {matiere}.

            Votre tâche est de générer une question de réflexion originale, en vous basant strictement sur les extraits suivants:

            {context}

            La question doit :
            1. Être ORIGINALE et non une reformulation des questions existantes
            2. Solliciter l'analyse critique d'un concept
            3. Être formulée de manière claire et précise
            4. Favoriser une réponse argumentée

            Instructions:
            - N'utilisez que les passages extraits des documents de cours ci-dessus
            - Ne proposez qu'une question ouverte et directe en une ligne
            - N'ajoutez ni explication ni sous-questions

            Votre question: 
            """
        
        return ChatPromptTemplate.from_template(template).partial(matiere=matiere)
    
    def _setup_rag_chain(
        self, 
        matiere: str, 
        custom_prompt: Optional[ChatPromptTemplate] = None,
        output_format: str = "text"
    ):
        """Configure la chaîne RAG pour une matière spécifique."""
        namespace = self._get_namespace(matiere)
        
        # Créer un store vectoriel pour cette matière
        vector_store = PineconeVectorStore(
            index_name=self.index_name,
            embedding=self.embeddings,
            namespace=namespace
        )
        
        # Configurer le prompt
        if custom_prompt:
            prompt = custom_prompt
        elif output_format == "json":
            prompt = self._create_json_prompt(matiere)
        else:
            prompt = hub.pull("langchain-ai/retrieval-qa-chat")
        
        # Configurer le retriever
        retriever = vector_store.as_retriever()
        
        # Configurer le modèle de langage
        llm = ChatOpenAI(
            openai_api_key=self.openai_api_key,
            model_name=self.llm_model,
            temperature=0.0
        )
        
        # Créer la chaîne de traitement des documents
        combine_docs_chain = create_stuff_documents_chain(llm, prompt)
        
        # Créer la chaîne de récupération
        retrieval_chain = create_retrieval_chain(retriever, combine_docs_chain)
        
        return retrieval_chain
    
    def query(
        self, 
        matiere: str, 
        question: str, 
        output_format: str = "text",
        custom_prompt: Optional[ChatPromptTemplate] = None
    ) -> Dict[str, Any]:
        """
        Interroge le système RAG pour une matière spécifique.
        
        Args:
            matiere: Identifiant de la matière
            question: Question à poser
            output_format: Format de sortie ("text" ou "json")
            custom_prompt: Prompt personnalisé optionnel
            
        Returns:
            Dictionnaire contenant la réponse et les métadonnées
        """
        # Configurer le système RAG
        retrieval_chain = self._setup_rag_chain(matiere, custom_prompt, output_format)
        
        print(f"Requête pour {matiere}: '{question}'")
        
        # Exécuter la requête
        response = retrieval_chain.invoke({"input": question})
        
        # Traiter la réponse selon le format
        if output_format == "json":
            try:
                # Extraire le JSON de la réponse
                json_answer = response['answer']
                
                # Rechercher le JSON entre les délimiteurs
                json_pattern = r"```(?:json)?(.*?)```"
                match = re.search(json_pattern, json_answer, re.DOTALL)
                
                if match:
                    json_str = match.group(1).strip()
                    response_json = json.loads(json_str)
                else:
                    json_str = json_answer.strip()
                    response_json = json.loads(json_str)
                
                # Ajouter les sources
                sources = []
                for i, doc in enumerate(response["context"]):
                    source_entry = {
                        "document": i + 1,
                        "source": doc.metadata.get('source', 'Source inconnue'),
                        "is_exam": doc.metadata.get('is_exam', False),
                        "contenu": doc.page_content[:250] + "..." if len(doc.page_content) > 250 else doc.page_content
                    }
                    sources.append(source_entry)
                
                # Enrichir la réponse JSON
                response_json.update({
                    "sources": sources,
                    "requête_originale": question,
                    "matière": matiere,
                    "date_génération": datetime.now().isoformat()
                })
                
                response['answer'] = json.dumps(response_json, ensure_ascii=False, indent=2)
                
            except Exception as e:
                print(f"Erreur lors du traitement JSON: {e}")
        
        return response
    
    def generate_reflection_question(
        self, 
        matiere: str, 
        concept_cle: Optional[str] = None,
        output_format: str = "text"
    ) -> Dict[str, Any]:
        """
        Génère une question de réflexion pour une matière.
        
        Args:
            matiere: Identifiant de la matière
            concept_cle: Concept spécifique (optionnel)
            output_format: Format de sortie ("text" ou "json")
            
        Returns:
            Dictionnaire contenant la question générée
        """
        # Créer le prompt tuteur
        tutor_prompt = self._create_tutor_prompt(matiere, output_format)
        
        # Construire la requête
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        if concept_cle and concept_cle.strip():
            query = f"Générer une question de réflexion sur le concept: {concept_cle} [timestamp: {timestamp}]"
        else:
            query = f"Générer une question de réflexion générale sur la matière [timestamp: {timestamp}]"
        
        # Interroger avec le prompt tuteur
        return self.query(matiere, query, output_format, tutor_prompt)
    
    def evaluate_student_response(
        self, 
        matiere: str, 
        question: str, 
        student_response: str
    ) -> Dict[str, Any]:
        """
        Évalue la réponse d'un étudiant à une question.
        
        Args:
            matiere: Identifiant de la matière
            question: Question posée
            student_response: Réponse de l'étudiant
            
        Returns:
            Dictionnaire contenant l'évaluation
        """
        evaluation_prompt = ChatPromptTemplate.from_template("""
        Vous êtes un examinateur académique spécialisé dans la matière {matiere}.
        
        Évaluez la réponse suivante en vous basant sur le contenu du cours:
        {context}
        
        Question: {question}
        Réponse de l'étudiant: {student_response}
        
        Votre évaluation DOIT être au format JSON suivant:
        
        ```json
        {{
            "note": 85,
            "commentaires": "Analyse détaillée de la réponse",
            "points_forts": ["Point fort 1", "Point fort 2"],
            "points_amélioration": ["Amélioration 1", "Amélioration 2"],
            "concepts_manqués": ["Concept 1", "Concept 2"]
        }}
        ```
        
        Ne répondez qu'avec ce format JSON.
        """).partial(matiere=matiere)
        
        # Construire la requête d'évaluation
        eval_query = f"Évaluer cette réponse: Question: {question} | Réponse: {student_response}"
        
        return self.query(matiere, eval_query, "json", evaluation_prompt)
    
    def get_index_stats(self) -> Dict[str, Any]:
        """Retourne les statistiques de l'index."""
        try:
            stats = self.index.describe_index_stats()
            return {
                "dimension": stats.dimension,
                "index_fullness": stats.index_fullness,
                "total_vector_count": stats.total_vector_count,
                "namespaces": {
                    ns_name: {"vector_count": ns_data.vector_count}
                    for ns_name, ns_data in (stats.namespaces or {}).items()
                }
            }
        except Exception as e:
            return {"error": str(e)}
    
    def list_matieres(self) -> List[str]:
        """Liste toutes les matières disponibles dans l'index."""
        try:
            stats = self.index.describe_index_stats()
            if stats.namespaces:
                return [
                    ns_name.replace("matiere-", "").upper()
                    for ns_name in stats.namespaces.keys()
                    if ns_name.startswith("matiere-")
                ]
            return []
        except Exception as e:
            print(f"Erreur lors de la récupération des matières: {e}")
            return []
    
    def delete_matiere(self, matiere: str) -> bool:
        """
        Supprime tous les documents d'une matière.
        
        Args:
            matiere: Identifiant de la matière à supprimer
            
        Returns:
            True si la suppression a réussi, False sinon
        """
        namespace = self._get_namespace(matiere)
        
        try:
            self.index.delete(delete_all=True, namespace=namespace)
            print(f"Matière {matiere} supprimée avec succès")
            return True
        except Exception as e:
            print(f"Erreur lors de la suppression de {matiere}: {e}")
            return False

    # Méthodes d'intégration pour l'architecture du projet
    
    def process_uploaded_file(self, matiere: str, file_path: str, is_exam: bool = False) -> Dict[str, Any]:
        """
        Traite un fichier uploadé et l'ajoute à l'index.
        
        Args:
            matiere: Identifiant de la matière
            file_path: Chemin vers le fichier uploadé
            is_exam: Si le fichier est un examen
            
        Returns:
            Dictionnaire avec les résultats du traitement
        """
        try:
            content = self._extract_content_from_file(file_path)
            if not content or content.strip() == "":
                return {
                    "success": False,
                    "message": "Impossible d'extraire le contenu du fichier",
                    "chunks_processed": 0
                }
            
            # Métadonnées du document
            filename = os.path.basename(file_path)
            metadata = {
                "source": filename,
                "matiere": matiere,
                "filename": filename,
                "filetype": os.path.splitext(filename)[1].lower(),
                "file_hash": self._calculate_file_hash(content),
                "updated_at": datetime.now().isoformat(),
                "is_exam": is_exam
            }
            
            # Découper le document
            documents = self._split_document(content, metadata)
            
            # Ajouter à l'index
            success = self.add_documents(matiere, documents)
            
            return {
                "success": success,
                "message": f"Document {filename} traité avec succès" if success else "Erreur lors du traitement",
                "chunks_processed": len(documents) if success else 0,
                "file_hash": metadata["file_hash"]
            }
            
        except Exception as e:
            return {
                "success": False,
                "message": f"Erreur lors du traitement: {str(e)}",
                "chunks_processed": 0
            }
    
    def get_documents_for_matiere(self, matiere: str) -> List[Dict[str, Any]]:
        """
        Récupère la liste des documents pour une matière.
        
        Args:
            matiere: Identifiant de la matière
            
        Returns:
            Liste des documents avec leurs métadonnées
        """
        try:
            namespace = self._get_namespace(matiere)
            
            # Pour l'instant, on retourne les statistiques du namespace
            # Dans une implémentation complète, on pourrait stocker les métadonnées
            # des documents dans une base de données séparée
            stats = self.index.describe_index_stats()
            
            if stats.namespaces and namespace in stats.namespaces:
                vector_count = stats.namespaces[namespace].vector_count
                return [{
                    "id": f"doc_{i}",
                    "filename": f"document_{i}.pdf",
                    "matiere": matiere,
                    "document_type": "pdf",
                    "is_exam": False,
                    "file_path": f"cours/{matiere}/document_{i}.pdf",
                    "file_size": 1024,
                    "upload_date": datetime.now().isoformat(),
                    "last_indexed": datetime.now().isoformat()
                } for i in range(min(vector_count // 10, 10))]  # Estimation
            
            return []
            
        except Exception as e:
            print(f"Erreur lors de la récupération des documents: {e}")
            return []


# Instance globale du service RAG (sera initialisée au démarrage de l'app)
rag_service: Optional[RAGService] = None


def get_rag_service() -> RAGService:
    """Retourne l'instance globale du service RAG."""
    global rag_service
    if rag_service is None:
        raise RuntimeError("Le service RAG n'est pas initialisé")
    return rag_service


def initialize_rag_service(
    pinecone_api_key: str,
    openai_api_key: str,
    index_name: str = "rhino-rag-index"
) -> bool:
    """
    Initialise le service RAG global.
    
    Args:
        pinecone_api_key: Clé API Pinecone
        openai_api_key: Clé API OpenAI
        index_name: Nom de l'index Pinecone
        
    Returns:
        True si l'initialisation a réussi, False sinon
    """
    global rag_service
    
    try:
        rag_service = RAGService(
            pinecone_api_key=pinecone_api_key,
            openai_api_key=openai_api_key,
            index_name=index_name
        )
        print("Service RAG initialisé avec succès")
        return True
    except Exception as e:
        print(f"Erreur lors de l'initialisation du service RAG: {e}")
        return False


# Exemple d'utilisation
if __name__ == "__main__":
    # Configuration (à adapter selon vos besoins)
    PINECONE_API_KEY = os.getenv("PINECONE_API_KEY", "your_pinecone_key")
    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "your_openai_key")
    
    # Initialiser le service RAG
    rag = RAGService(
        pinecone_api_key=PINECONE_API_KEY,
        openai_api_key=OPENAI_API_KEY,
        index_name="rag-example"
    )
    
    # Exemple d'ajout de documents depuis un répertoire
    # rag.add_documents_from_directory("SYD", "/path/to/syd/documents")
    
    # Exemple de requête
    # response = rag.query("SYD", "Qu'est-ce que TCP/IP?", output_format="json")
    # print(response['answer'])
    
    # Exemple de génération de question
    # question = rag.generate_reflection_question("SYD", "protocoles réseau", output_format="json")
    # print(question['answer'])
    
    # Afficher les statistiques
    print("Statistiques de l'index:")
    print(json.dumps(rag.get_index_stats(), indent=2))
    
    print("Matières disponibles:")
    print(rag.list_matieres()) 